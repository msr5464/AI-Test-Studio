"""
Multi-Format RAG System
=======================

Unified RAG system that supports multiple document formats:
- Text strings
- PDF files
- CSV files
- Excel files (.xlsx, .xls)
- Word documents (.docx, .doc)
- PowerPoint presentations (.pptx, .ppt)

Automatically detects file format and uses appropriate loader.
All format-specific logic is consolidated in this single file.
"""

import sys
import os
import re
from pathlib import Path
from typing import List, Union, Optional

# Handle both relative and absolute imports
try:
    from .base_rag import BaseRAG
    from .rag_helpers import calculate_file_hash, add_file_metadata_to_documents
    from langchain_core.documents import Document
except ImportError:
    from backend.rag.base_rag import BaseRAG
    from backend.rag.rag_helpers import calculate_file_hash, add_file_metadata_to_documents
    from langchain_core.documents import Document

# PDF loader
try:
    from langchain_community.document_loaders import PyPDFLoader
    PDF_LOADER_AVAILABLE = True
except ImportError:
    PDF_LOADER_AVAILABLE = False

# CSV/Excel loaders
try:
    from langchain_community.document_loaders import CSVLoader, UnstructuredExcelLoader
    CSV_EXCEL_LOADER_AVAILABLE = True
except ImportError:
    CSV_EXCEL_LOADER_AVAILABLE = False

# Pandas for CSV/Excel processing
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# Word document loader
try:
    from docx import Document as DocxDocument
    DOCX_LOADER_AVAILABLE = True
except ImportError:
    DOCX_LOADER_AVAILABLE = False

# Unstructured Word loader (supports both .doc and .docx)
try:
    from langchain_community.document_loaders import UnstructuredWordDocumentLoader
    UNSTRUCTURED_WORD_LOADER_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_WORD_LOADER_AVAILABLE = False

# Check if any Word loader is available
WORD_LOADER_AVAILABLE = DOCX_LOADER_AVAILABLE or UNSTRUCTURED_WORD_LOADER_AVAILABLE

# PowerPoint loader
try:
    from pptx import Presentation
    PPTX_LOADER_AVAILABLE = True
except ImportError:
    PPTX_LOADER_AVAILABLE = False

# Unstructured PowerPoint loader (supports both .ppt and .pptx)
try:
    from langchain_community.document_loaders import UnstructuredPowerPointLoader
    UNSTRUCTURED_PPT_LOADER_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_PPT_LOADER_AVAILABLE = False

# Check if any PowerPoint loader is available
POWERPOINT_LOADER_AVAILABLE = PPTX_LOADER_AVAILABLE or UNSTRUCTURED_PPT_LOADER_AVAILABLE


class MultiFormatRAG(BaseRAG):
    """
    Unified RAG system supporting multiple document formats.
    
    Supports:
    - Text strings (in-memory)
    - PDF files (.pdf)
    - CSV files (.csv)
    - Excel files (.xlsx, .xls)
    - Word documents (.docx, .doc)
    - PowerPoint presentations (.pptx, .ppt)
    
    Automatically detects format and uses appropriate processing.
    All format-specific logic is consolidated in this class.
    """
    
    def __init__(self, *args, csv_separator: str = ",", excel_sheet_names: Optional[List[str]] = None,
                 config: Optional[object] = None, **kwargs):
        """
        Initialize multi-format RAG system.
        
        Args:
            *args: Arguments passed to BaseRAG
            csv_separator: CSV separator character. Default: ","
            excel_sheet_names: List of sheet names to load (None = all sheets). Default: None
            config: Optional RAGConfig object. If provided, overrides individual parameters.
            **kwargs: Keyword arguments passed to BaseRAG
        """
        # If config is provided, extract its values
        if config:
            try:
                # Try relative import first (when used as module)
                try:
                    from .settings import RAGConfig
                except ImportError:
                    # Fallback to absolute import (when run as script)
                    from settings import RAGConfig
                if isinstance(config, RAGConfig):
                    # Use to_dict() which excludes CSV/Excel specific settings
                    kwargs.update(config.to_dict(exclude_testrail=True))
            except ImportError:
                pass
        
        super().__init__(*args, config=config, **kwargs)
        
        # CSV/Excel specific settings
        self.csv_separator = csv_separator
        self.excel_sheet_names = excel_sheet_names
    
    def _clean_pdf_text(self, text: str) -> str:
        """
        Clean PDF text by removing excessive spaces and normalizing whitespace.
        
        Args:
            text: Raw text from PDF
            
        Returns:
            Cleaned text
        """
        # Remove multiple spaces and normalize whitespace
        cleaned = ' '.join(text.split())
        # Fix common PDF extraction issues: single letters separated by spaces
        # Pattern: "A s p i r e" -> "Aspire"
        cleaned = re.sub(r'\b([A-Za-z])\s+(?=[A-Za-z]\s+[A-Za-z])', r'\1', cleaned)
        # Fix patterns like "Q A" -> "QA", but preserve "Q A -" -> "QA -"
        cleaned = re.sub(r'\b([A-Z])\s+([A-Z])\b(?!\s*-)', r'\1\2', cleaned)
        # Normalize multiple spaces again after fixes
        cleaned = ' '.join(cleaned.split())
        return cleaned
    
    def _convert_dataframe_to_text(self, df: 'pd.DataFrame', sheet_name: Optional[str] = None) -> str:
        """
        Convert pandas DataFrame to text format for RAG.
        
        Args:
            df: Pandas DataFrame
            sheet_name: Optional sheet name (for Excel files)
            
        Returns:
            Text representation of the DataFrame
        """
        if df.empty:
            return ""
        
        lines = []
        
        # Add sheet name if provided
        if sheet_name:
            lines.append(f"Sheet: {sheet_name}")
            lines.append("=" * 60)
        
        # Add column headers
        headers = " | ".join(str(col) for col in df.columns)
        lines.append(f"Columns: {headers}")
        lines.append("-" * 60)
        
        # Convert each row to text
        for idx, row in df.iterrows():
            row_text_parts = []
            for col in df.columns:
                value = str(row[col]) if pd.notna(row[col]) else "N/A"
                row_text_parts.append(f"{col}: {value}")
            lines.append(f"Row {idx + 1}: {' | '.join(row_text_parts)}")
        
        return "\n".join(lines)

    def _clean_html(self, text: str) -> str:
        """
        Clean HTML tags from text while preserving meaningful whitespace and alignment.
        """
        if not isinstance(text, str) or not text or text.lower() == "nan":
            return ""
            
        import re
        import html
        
        # 1. Replace block-level elements with newlines to maintain structure
        block_tags = [
            'p', 'div', 'li', 'br', 'ul', 'ol', 'tr', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
            'table', 'thead', 'tbody', 'tfoot', 'caption', 'blockquote', 'hr'
        ]
        tag_pattern = r'</?(?:' + '|'.join(block_tags) + r')[^>]*>'
        text = re.sub(tag_pattern, '\n', text, flags=re.IGNORECASE)
        
        # 2. Strip all remaining HTML tags
        text = re.sub(r'<[^>]+>', '', text)
        
        # 3. Handle HTML entities
        text = html.unescape(text)
        
        # 4. Process line by line
        lines = []
        for line in text.split('\n'):
            # Strip whitespace and normalize trailing backslashes (often found in shell commands in HTML/CSV)
            cleaned_line = line.strip().rstrip('\\').strip()
            if cleaned_line:
                # Clean up double quotes common in CSV exports of JSON/String data
                cleaned_line = cleaned_line.replace('""', '"')
                lines.append(cleaned_line)
            
        # 5. Join with newlines and collapse excessive vertical space
        result = '\n'.join(lines)
        result = re.sub(r'\n{3,}', '\n\n', result)
        
        return result.strip()

    def _is_testcase_csv(self, df: 'pd.DataFrame') -> bool:
        """Detect if a DataFrame contains testcase schema based on headers."""
        if df.empty:
            return False
            
        target_headers = {
            'id', 'title', 'execution mode', 'expected result', 
            'platform', 'preconditions', 'priority', 
            'section hierarchy', 'steps', 'type'
        }
        
        columns = {str(col).lower().strip() for col in df.columns}
        matched = columns.intersection(target_headers)
        
        # If at least 70% match (7 out of 10), consider it a testcase CSV
        return len(matched) >= 7

    def _is_specs_csv(self, df: 'pd.DataFrame') -> bool:
        """Check if DataFrame has specs structure: page_id, title, body (e.g. Confluence sync)."""
        target_headers = {'page_id', 'title', 'body'}
        columns = {str(col).lower().strip() for col in df.columns}
        return target_headers.issubset(columns)

    def _convert_specs_dataframe_to_documents(self, df: 'pd.DataFrame') -> List[Document]:
        """Convert specs DataFrame (page_id, title, body) to Documents with label 'specs'."""
        documents = []
        cols = {str(col).lower().strip(): col for col in df.columns}

        for idx, row in df.iterrows():
            def get_clean(header_key):
                if header_key in cols:
                    val = row[cols[header_key]]
                    return self._clean_html(str(val)) if pd.notna(val) else ""
                return ""

            page_id = get_clean('page_id') or f"Page_{idx + 1}"
            title = get_clean('title') or "Untitled"
            body = get_clean('body') or ""
            url = get_clean('url') if 'url' in cols else ""

            content = f"Spec (Page ID: {page_id})\nTitle: {title}\n"
            if url:
                content += f"URL: {url}\n"
            content += f"\nContent:\n{body}"

            metadata = {
                "page_id": page_id,
                "label": "specs",
                "source_type": "specs",
            }
            if url:
                metadata["url"] = url

            documents.append(Document(page_content=content.strip(), metadata=metadata))

        return documents

    def _convert_testcase_dataframe_to_documents(self, df: 'pd.DataFrame') -> List[Document]:
        """Specialized conversion for testcase dataframes."""
        documents = []
        
        # Get column index mapping (to handle casing differences)
        cols = {str(col).lower().strip(): col for col in df.columns}
        
        for idx, row in df.iterrows():
            # Extract and clean mapped fields
            def get_clean(header_key):
                if header_key in cols:
                    val = row[cols[header_key]]
                    return self._clean_html(str(val)) if pd.notna(val) else ""
                return ""

            tc_id = get_clean('id') or f"Row_{idx+1}"
            tc_priority = get_clean('priority') or "N/A"
            tc_platform = get_clean('platform') or "N/A"
            tc_type = get_clean('type') or ""   # e.g. "Functional", "FCT / Regression", "Prod Sanity"
            tc_title = get_clean('title') or "Untitled"
            
            section = get_clean('section hierarchy')
            preconditions = get_clean('preconditions')
            steps = get_clean('steps')
            expected = get_clean('expected result')
            requirement = get_clean('requirement')
            
            # Format page content with clear sections (for retrieval and backward compatibility).
            # When Requirement column exists (e.g. from pushed generated tests), put it first
            # so retrieval by requirement text matches this test.
            content = ""
            if requirement:
                content += f"Requirement: {requirement}\n\n"
            content += f"Testrail Id: {tc_id}\n"
            content += f"Priority: {tc_priority}\n"
            content += f"Platform: {tc_platform}\n"
            if tc_type:
                content += f"Type: {tc_type}\n"
            content += f"Title: {tc_title}\n"
            if section:
                content += f"Section Hierarchy: {section}\n"
            if preconditions:
                content += f"Preconditions: {preconditions}\n"
            if steps:
                content += f"Steps: {steps}\n"
            if expected:
                content += f"Expected Result: {expected}\n"
            
            # Add other fields if they exist but aren't in the main mapping
            other_parts = []
            for raw_col in df.columns:
                lower_col = str(raw_col).lower().strip()
                if lower_col not in {'id', 'priority', 'platform', 'title', 'section hierarchy', 'preconditions', 'steps', 'expected result', 'requirement', 'type'}:
                    val = row[raw_col]
                    cleaned_val = self._clean_html(str(val)) if pd.notna(val) else ""
                    if cleaned_val:
                        other_parts.append(f"{raw_col}: {cleaned_val}")
            
            if other_parts:
                content += "\n".join(other_parts)
            
            # Store Preconditions, Steps, Expected Result separately in metadata
            metadata = {
                "testrail_id": tc_id,
                "priority": tc_priority,
                "platform": tc_platform,
                "case_type": tc_type,
                "source_type": "testcase",
                "preconditions": preconditions or "",
                "steps": steps or "",
                "expected_result": expected or "",
            }
            
            documents.append(Document(page_content=content.strip(), metadata=metadata))
            
        return documents
    
    def _load_csv(self, csv_path: Union[str, Path]) -> List[Document]:
        """
        Load CSV file and convert to text documents.
        
        Args:
            csv_path: Path to CSV file
            
        Returns:
            List of text documents (one per chunk)
        """
        csv_path = Path(csv_path)
        
        if not csv_path.exists():
            raise FileNotFoundError(f"CSV file not found: {csv_path}")
        
        if csv_path.suffix.lower() != '.csv':
            raise ValueError(f"File must be a CSV: {csv_path}")
        
        print(f"📊 Loading CSV: {csv_path.name}")
        
        if PANDAS_AVAILABLE:
            # Use pandas for better control
            try:
                df = pd.read_csv(csv_path, sep=self.csv_separator)
                print(f"   Loaded {len(df)} rows, {len(df.columns)} columns")
                
                # Detect if it's a specs CSV (page_id, title, body - e.g. Confluence sync)
                if self._is_specs_csv(df):
                    print(f"📋 Detected Specs CSV: Using label 'specs'")
                    return self._convert_specs_dataframe_to_documents(df)
                # Detect if it's a testcase CSV
                if self._is_testcase_csv(df):
                    print(f"📋 Detected Testcase CSV: Using specialized mapping and cleaning")
                    return self._convert_testcase_dataframe_to_documents(df)
                
                # Default behavior: clean every cell and convert to text
                df_cleaned = df.copy()
                for col in df.columns:
                    df_cleaned[col] = df_cleaned[col].apply(lambda x: self._clean_html(str(x)) if pd.notna(x) else "")
                
                text = self._convert_dataframe_to_text(df_cleaned)
                
                # Split into chunks if needed
                doc = Document(page_content=text)
                return self.text_splitter.split_documents([doc])
            except Exception as e:
                print(f"⚠️  Error loading CSV with pandas: {e}")
                # Fallback to CSVLoader
                if CSV_EXCEL_LOADER_AVAILABLE:
                    loader = CSVLoader(str(csv_path))
                    return loader.load()
                else:
                    raise
        elif CSV_EXCEL_LOADER_AVAILABLE:
            # Fallback to LangChain CSVLoader
            loader = CSVLoader(str(csv_path))
            return loader.load()
        else:
            raise ImportError("Neither pandas nor CSVLoader available. Install: pip install pandas")
    
    def _load_excel(self, excel_path: Union[str, Path]) -> List[str]:
        """
        Load Excel file and convert to text documents.
        
        Args:
            excel_path: Path to Excel file
            
        Returns:
            List of text documents (one per sheet/chunk)
        """
        excel_path = Path(excel_path)
        
        if not excel_path.exists():
            raise FileNotFoundError(f"Excel file not found: {excel_path}")
        
        valid_extensions = ['.xlsx', '.xls']
        if excel_path.suffix.lower() not in valid_extensions:
            raise ValueError(f"File must be an Excel file (.xlsx or .xls): {excel_path}")
        
        print(f"📊 Loading Excel: {excel_path.name}")
        
        if PANDAS_AVAILABLE:
            # Use pandas for better control
            try:
                all_texts = []
                
                # Read all sheets or specified sheets
                if self.excel_sheet_names:
                    sheet_names = self.excel_sheet_names
                else:
                    # Read all sheets
                    excel_file = pd.ExcelFile(excel_path)
                    sheet_names = excel_file.sheet_names
                
                print(f"   Found {len(sheet_names)} sheet(s)")
                
                for sheet_name in sheet_names:
                    try:
                        df = pd.read_excel(excel_path, sheet_name=sheet_name)
                        print(f"   Loaded sheet '{sheet_name}': {len(df)} rows, {len(df.columns)} columns")
                        
                        # Convert to text
                        text = self._convert_dataframe_to_text(df, sheet_name=sheet_name)
                        all_texts.append(text)
                    except Exception as e:
                        print(f"⚠️  Error loading sheet '{sheet_name}': {e}")
                        continue
                
                # Split into chunks
                documents = []
                for text in all_texts:
                    doc = Document(page_content=text)
                    splits = self.text_splitter.split_documents([doc])
                    documents.extend([split.page_content for split in splits])
                
                return documents
            except Exception as e:
                print(f"⚠️  Error loading Excel with pandas: {e}")
                # Fallback to UnstructuredExcelLoader
                if CSV_EXCEL_LOADER_AVAILABLE:
                    loader = UnstructuredExcelLoader(str(excel_path))
                    documents = loader.load()
                    return [doc.page_content for doc in documents]
                else:
                    raise
        elif CSV_EXCEL_LOADER_AVAILABLE:
            # Fallback to LangChain UnstructuredExcelLoader
            loader = UnstructuredExcelLoader(str(excel_path))
            documents = loader.load()
            return [doc.page_content for doc in documents]
        else:
            raise ImportError("Neither pandas nor UnstructuredExcelLoader available. Install: pip install pandas openpyxl")
    
    def _load_pdf(self, pdf_path: Union[str, Path]) -> List[Document]:
        """
        Load PDF file and return documents.
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of Document objects
        """
        pdf_path = Path(pdf_path)
        
        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        if not pdf_path.suffix.lower() == '.pdf':
            raise ValueError(f"File must be a PDF: {pdf_path}")
        
        if not PDF_LOADER_AVAILABLE:
            raise ImportError("PyPDFLoader required. Install: pip install pypdf")
        
        print(f"📄 Loading PDF: {pdf_path.name}")
        
        loader = PyPDFLoader(str(pdf_path))
        documents = loader.load()
        print(f"   Loaded {len(documents)} pages")
        
        # Clean PDF text - remove excessive spaces and normalize whitespace
        for doc in documents:
            doc.page_content = self._clean_pdf_text(doc.page_content)
        
        return documents
    
    def _load_word(self, word_path: Union[str, Path]) -> List[Document]:
        """
        Load Word document and return documents.
        
        Args:
            word_path: Path to Word document (.docx or .doc)
            
        Returns:
            List of Document objects
        """
        word_path = Path(word_path)
        
        if not word_path.exists():
            raise FileNotFoundError(f"Word document not found: {word_path}")
        
        valid_extensions = ['.docx', '.doc', '.docs']
        if word_path.suffix.lower() not in valid_extensions:
            raise ValueError(f"File must be a Word document (.docx, .doc, or .docs): {word_path}")
        
        if not WORD_LOADER_AVAILABLE:
            raise ImportError("Word document support requires either python-docx or unstructured. Install: pip install python-docx or pip install unstructured")
        
        print(f"📄 Loading Word document: {word_path.name}")
        
        file_ext = word_path.suffix.lower()
        
        # Use python-docx for .docx files (faster and more reliable)
        if file_ext == '.docx' and DOCX_LOADER_AVAILABLE:
            try:
                docx = DocxDocument(str(word_path))
                paragraphs = []
                
                # Extract text from all paragraphs
                for para in docx.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text.strip())
                
                # Extract text from tables
                for table in docx.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            paragraphs.append(" | ".join(row_text))
                
                # Combine all paragraphs into a single document
                full_text = "\n\n".join(paragraphs)
                
                if not full_text.strip():
                    print(f"⚠️  Warning: No text content found in {word_path.name}")
                    return []
                
                print(f"   Extracted {len(paragraphs)} paragraph(s)")
                
                # Create a single document
                return [Document(page_content=full_text)]
                
            except Exception as e:
                print(f"⚠️  Error loading .docx with python-docx: {e}")
                # Fall back to UnstructuredWordDocumentLoader if available
                if UNSTRUCTURED_WORD_LOADER_AVAILABLE:
                    print(f"   Trying UnstructuredWordDocumentLoader as fallback...")
                else:
                    raise
        
        # Use UnstructuredWordDocumentLoader for .doc and .docs files (or as fallback for .docx)
        if file_ext in ['.doc', '.docs'] or (file_ext == '.docx' and not DOCX_LOADER_AVAILABLE):
            if not UNSTRUCTURED_WORD_LOADER_AVAILABLE:
                raise ImportError("UnstructuredWordDocumentLoader required for .doc/.docs files. Install: pip install unstructured")
            
            # Check if python-docx is available (required by UnstructuredWordDocumentLoader)
            if not DOCX_LOADER_AVAILABLE:
                raise ImportError(
                    "python-docx is required for Word document processing. "
                    "UnstructuredWordDocumentLoader depends on it. "
                    "Install: pip install python-docx"
                )
            
            try:
                loader = UnstructuredWordDocumentLoader(str(word_path))
                documents = loader.load()
                
                if not documents:
                    print(f"⚠️  Warning: No text content found in {word_path.name}")
                    return []
                
                # Combine all documents into a single document
                full_text = "\n\n".join([doc.page_content for doc in documents if doc.page_content.strip()])
                
                if not full_text.strip():
                    print(f"⚠️  Warning: No text content found in {word_path.name}")
                    return []
                
                print(f"   Loaded {len(documents)} document section(s)")
                
                return [Document(page_content=full_text)]
                
            except ImportError as e:
                if 'docx' in str(e).lower():
                    raise ImportError(
                        f"Missing dependency for Word document processing: {e}\n"
                        "Please install: pip install python-docx"
                    )
                raise
            except Exception as e:
                error_msg = str(e)
                # Check if it's a LibreOffice error
                if 'soffice' in error_msg.lower() or 'libreoffice' in error_msg.lower():
                    raise ImportError(
                        f"LibreOffice is required for .doc file processing.\n\n"
                        f"Error: {error_msg}\n\n"
                        f"Installation instructions:\n"
                        f"  - macOS: brew install --cask libreoffice\n"
                        f"  - Linux (Debian/Ubuntu): sudo apt-get install libreoffice\n"
                        f"  - Linux (RHEL/CentOS): sudo yum install libreoffice\n"
                        f"  - Windows: Download from https://www.libreoffice.org/download/\n\n"
                        f"Alternatively, convert .doc files to .docx format before uploading."
                    )
                print(f"⚠️  Error loading Word document with UnstructuredWordDocumentLoader: {e}")
                raise
    
    def _load_powerpoint(self, pptx_path: Union[str, Path]) -> List[Document]:
        """
        Load PowerPoint presentation and return documents.
        
        Args:
            pptx_path: Path to PowerPoint file (.pptx or .ppt)
            
        Returns:
            List of Document objects (one per slide)
        """
        pptx_path = Path(pptx_path)
        
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        valid_extensions = ['.pptx', '.ppt']
        if pptx_path.suffix.lower() not in valid_extensions:
            raise ValueError(f"File must be a PowerPoint file (.pptx or .ppt): {pptx_path}")
        
        if not POWERPOINT_LOADER_AVAILABLE:
            raise ImportError("PowerPoint support requires either python-pptx or unstructured. Install: pip install python-pptx or pip install unstructured")
        
        print(f"📊 Loading PowerPoint: {pptx_path.name}")
        
        file_ext = pptx_path.suffix.lower()
        
        # Use python-pptx for .pptx files (faster and more reliable)
        if file_ext == '.pptx' and PPTX_LOADER_AVAILABLE:
            try:
                prs = Presentation(str(pptx_path))
                documents = []
                
                # Extract text from each slide
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    
                    # Extract text from all shapes on the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    
                    # Combine all text from the slide
                    if slide_texts:
                        slide_content = "\n".join(slide_texts)
                        # Create a document for each slide with metadata
                        doc = Document(
                            page_content=slide_content,
                            metadata={"slide_number": slide_num, "total_slides": len(prs.slides)}
                        )
                        documents.append(doc)
                
                if not documents:
                    print(f"⚠️  Warning: No text content found in {pptx_path.name}")
                    return []
                
                print(f"   Extracted {len(documents)} slide(s)")
                
                return documents
                
            except Exception as e:
                print(f"⚠️  Error loading .pptx with python-pptx: {e}")
                # Fall back to UnstructuredPowerPointLoader if available
                if UNSTRUCTURED_PPT_LOADER_AVAILABLE:
                    print(f"   Trying UnstructuredPowerPointLoader as fallback...")
                else:
                    raise
        
        # Use UnstructuredPowerPointLoader for .ppt files (or as fallback for .pptx)
        if file_ext == '.ppt' or (file_ext == '.pptx' and not PPTX_LOADER_AVAILABLE):
            if not UNSTRUCTURED_PPT_LOADER_AVAILABLE:
                raise ImportError("UnstructuredPowerPointLoader required for .ppt files. Install: pip install unstructured")
            
            try:
                loader = UnstructuredPowerPointLoader(str(pptx_path))
                documents = loader.load()
                
                if not documents:
                    print(f"⚠️  Warning: No text content found in {pptx_path.name}")
                    return []
                
                print(f"   Loaded {len(documents)} document section(s)")
                
                return documents
                
            except ImportError as e:
                if 'pptx' in str(e).lower() or 'ppt' in str(e).lower():
                    raise ImportError(
                        f"Missing dependency for PowerPoint processing: {e}\n"
                        "Please install: pip install python-pptx"
                    )
                raise
            except Exception as e:
                print(f"⚠️  Error loading PowerPoint file with UnstructuredPowerPointLoader: {e}")
                raise
    
    def _detect_format(self, file_path: Union[str, Path]) -> str:
        """
        Detect file format from extension.
        
        Args:
            file_path: Path to file
            
        Returns:
            Format string: 'pdf', 'csv', 'excel', 'word', 'powerpoint', 'text', or 'unknown'
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        if suffix == '.pdf':
            return 'pdf'
        elif suffix == '.csv':
            return 'csv'
        elif suffix in ['.xlsx', '.xls']:
            return 'excel'
        elif suffix in ['.docx', '.doc', '.docs']:
            return 'word'
        elif suffix in ['.pptx', '.ppt']:
            return 'powerpoint'
        elif suffix in ['.txt', '.text', '.md']:
            return 'text'
        else:
            return 'unknown'
    
    def _load_text_file(self, text_path: Union[str, Path]) -> str:
        """
        Load text file and return content.
        
        Args:
            text_path: Path to text file
            
        Returns:
            File content as string
        """
        text_path = Path(text_path)
        
        if not text_path.exists():
            raise FileNotFoundError(f"Text file not found: {text_path}")
        
        valid_extensions = ['.txt', '.text', '.md']
        if text_path.suffix.lower() not in valid_extensions:
            raise ValueError(f"File must be a text file (.txt, .text, or .md): {text_path}")
        
        print(f"📄 Loading text file: {text_path.name}")
        
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
            content = None
            
            for encoding in encodings:
                try:
                    with open(text_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise ValueError(f"Could not decode text file with any supported encoding: {text_path}")
            
            if not content.strip():
                print(f"⚠️  Warning: Text file is empty: {text_path.name}")
                return ""
            
            print(f"   Loaded {len(content)} characters")
            return content
            
        except Exception as e:
            print(f"⚠️  Error loading text file: {e}")
            raise
    
    def add_texts(self, texts: List[str]):
        """
        Add text documents from strings.
        
        Args:
            texts: List of text strings
        """
        if not texts:
            return
        
        print(f"📚 Adding {len(texts)} text document(s)...")
        
        # Use BaseRAG's method for text documents
        self.add_documents_from_texts(texts)
    
    def add_text_file(self, text_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add text file(s). Accepts either a single file path or a list of file paths.
        
        Args:
            text_paths: Single text file path or list of text file paths (.txt, .text, .md)
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(text_paths, list):
            text_paths = [text_paths]
        
        print(f"📚 Loading {len(text_paths)} text file(s)...")
        
        all_splits = []
        processed_files = []
        
        for text_path in text_paths:
            text_path = Path(text_path).absolute()
            if not text_path.exists():
                print(f"⚠️  Skipping non-existent file: {text_path}")
                continue
            
            valid_extensions = ['.txt', '.text', '.md']
            if text_path.suffix.lower() not in valid_extensions:
                print(f"⚠️  Skipping non-text file: {text_path}")
                continue
            
            try:
                content = self._load_text_file(text_path)
                
                if not content.strip():
                    print(f"⚠️  Skipping empty file: {text_path.name}")
                    continue
                
                # Create document from text content
                document = Document(page_content=content)
                
                # Split document into chunks
                splits = self.text_splitter.split_documents([document])
                
                # Add file metadata
                file_hash = calculate_file_hash(text_path)
                splits = add_file_metadata_to_documents(splits, file_path=text_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(text_path)
            except Exception as e:
                print(f"⚠️  Error loading {text_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid text files loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All text files indexed!"
        if removed_count > 0:
            success_msg = f"All text files indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_pdf(self, pdf_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add PDF file(s). Accepts either a single file path or a list of file paths.
        
        Args:
            pdf_paths: Single PDF file path or list of PDF file paths
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(pdf_paths, list):
            pdf_paths = [pdf_paths]
        
        print(f"📚 Loading {len(pdf_paths)} PDF file(s)...")
        
        all_splits = []
        processed_files = []
        
        for pdf_path in pdf_paths:
            pdf_path = Path(pdf_path).absolute()
            if not pdf_path.exists():
                print(f"⚠️  Skipping non-existent file: {pdf_path}")
                continue
            
            if not pdf_path.suffix.lower() == '.pdf':
                print(f"⚠️  Skipping non-PDF file: {pdf_path}")
                continue
            
            try:
                documents = self._load_pdf(pdf_path)
                
                # Split documents into chunks
                splits = self.text_splitter.split_documents(documents)
                
                # Add file metadata
                file_hash = calculate_file_hash(pdf_path)
                splits = add_file_metadata_to_documents(splits, file_path=pdf_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(pdf_path)
            except Exception as e:
                print(f"⚠️  Error loading {pdf_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid PDF files loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All PDFs indexed!"
        if removed_count > 0:
            success_msg = f"All PDFs indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_csv(self, csv_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add CSV file(s). Accepts either a single file path or a list of file paths.
        
        Args:
            csv_paths: Single CSV file path or list of CSV file paths
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(csv_paths, list):
            csv_paths = [csv_paths]
        
        print(f"📚 Loading {len(csv_paths)} CSV file(s)...")
        
        all_splits = []
        processed_files = []
        
        for csv_path in csv_paths:
            csv_path = Path(csv_path).absolute()
            if not csv_path.exists():
                print(f"⚠️  Skipping non-existent file: {csv_path}")
                continue
            
            if csv_path.suffix.lower() != '.csv':
                print(f"⚠️  Skipping non-CSV file: {csv_path}")
                continue
            
            try:
                splits = self._load_csv(csv_path)
                
                # Add file metadata
                file_hash = calculate_file_hash(csv_path)
                splits = add_file_metadata_to_documents(splits, file_path=csv_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(csv_path)
            except Exception as e:
                print(f"⚠️  Error loading {csv_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid CSV files loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All CSV files indexed!"
        if removed_count > 0:
            success_msg = f"All CSV files indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_excel(self, excel_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add Excel file(s). Accepts either a single file path or a list of file paths.
        
        Args:
            excel_paths: Single Excel file path or list of Excel file paths
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(excel_paths, list):
            excel_paths = [excel_paths]
        
        print(f"📚 Loading {len(excel_paths)} Excel file(s)...")
        
        all_splits = []
        processed_files = []
        
        for excel_path in excel_paths:
            excel_path = Path(excel_path).absolute()
            if not excel_path.exists():
                print(f"⚠️  Skipping non-existent file: {excel_path}")
                continue
            
            valid_extensions = ['.xlsx', '.xls']
            if excel_path.suffix.lower() not in valid_extensions:
                print(f"⚠️  Skipping non-Excel file: {excel_path}")
                continue
            
            try:
                texts = self._load_excel(excel_path)
                documents = [Document(page_content=text) for text in texts]
                splits = self.text_splitter.split_documents(documents)
                
                # Add file metadata
                file_hash = calculate_file_hash(excel_path)
                splits = add_file_metadata_to_documents(splits, file_path=excel_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(excel_path)
            except Exception as e:
                print(f"⚠️  Error loading {excel_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid Excel files loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All Excel files indexed!"
        if removed_count > 0:
            success_msg = f"All Excel files indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_word(self, word_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add Word document(s). Accepts either a single file path or a list of file paths.
        
        Args:
            word_paths: Single Word file path or list of Word file paths (.docx)
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(word_paths, list):
            word_paths = [word_paths]
        
        print(f"📚 Loading {len(word_paths)} Word document(s)...")
        
        all_splits = []
        processed_files = []
        
        for word_path in word_paths:
            word_path = Path(word_path).absolute()
            if not word_path.exists():
                print(f"⚠️  Skipping non-existent file: {word_path}")
                continue
            
            valid_extensions = ['.docx', '.doc', '.docs']
            if word_path.suffix.lower() not in valid_extensions:
                print(f"⚠️  Skipping non-Word file: {word_path}")
                continue
            
            try:
                documents = self._load_word(word_path)
                
                # Split documents into chunks
                splits = self.text_splitter.split_documents(documents)
                
                # Add file metadata
                file_hash = calculate_file_hash(word_path)
                splits = add_file_metadata_to_documents(splits, file_path=word_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(word_path)
            except Exception as e:
                print(f"⚠️  Error loading {word_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid Word documents loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All Word documents indexed!"
        if removed_count > 0:
            success_msg = f"All Word documents indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_powerpoint(self, pptx_paths: Union[str, Path, List[Union[str, Path]]], replace_if_exists: bool = True):
        """
        Add PowerPoint presentation(s). Accepts either a single file path or a list of file paths.
        
        Args:
            pptx_paths: Single PowerPoint file path or list of PowerPoint file paths (.pptx)
            replace_if_exists: If True, replace existing documents from the same files. Default: True
        """
        # Normalize to list
        if not isinstance(pptx_paths, list):
            pptx_paths = [pptx_paths]
        
        print(f"📚 Loading {len(pptx_paths)} PowerPoint file(s)...")
        
        all_splits = []
        processed_files = []
        
        for pptx_path in pptx_paths:
            pptx_path = Path(pptx_path).absolute()
            if not pptx_path.exists():
                print(f"⚠️  Skipping non-existent file: {pptx_path}")
                continue
            
            valid_extensions = ['.pptx', '.ppt']
            if pptx_path.suffix.lower() not in valid_extensions:
                print(f"⚠️  Skipping non-PowerPoint file: {pptx_path}")
                continue
            
            try:
                documents = self._load_powerpoint(pptx_path)
                
                # Split documents into chunks
                splits = self.text_splitter.split_documents(documents)
                
                # Add file metadata
                file_hash = calculate_file_hash(pptx_path)
                splits = add_file_metadata_to_documents(splits, file_path=pptx_path, file_hash=file_hash)
                
                all_splits.extend(splits)
                processed_files.append(pptx_path)
            except Exception as e:
                print(f"⚠️  Error loading {pptx_path.name}: {e}")
                continue
        
        print(f"📄 Total chunks: {len(all_splits)}")
        
        if not all_splits:
            print("⚠️  No valid PowerPoint files loaded")
            return
        
        # Load vectorstore if needed and remove duplicates
        self._load_vectorstore_if_needed()
        removed_count = self._remove_existing_documents(processed_files, replace_if_exists)
        
        # Add documents to vectorstore
        success_msg = f"All PowerPoint files indexed!"
        if removed_count > 0:
            success_msg = f"All PowerPoint files indexed! (Replaced {removed_count} existing document(s))"
        self._add_documents_to_vectorstore(all_splits, success_msg)
    
    def add_documents(self, inputs: Union[List[str], List[Union[str, Path]]]):
        """
        Add documents - automatically detects format.
        
        Args:
            inputs: Can be:
                - List of text strings (for text documents)
                - List of file paths (for PDF/CSV/Excel/Word/PowerPoint files)
        """
        if not inputs:
            return
        
        # Check if first item is a file path or text string
        first_item = inputs[0]
        
        # If it's a Path object or string that looks like a file path
        if isinstance(first_item, (str, Path)):
            path = Path(first_item)
            # Check if it exists and is a file
            if path.exists() and path.is_file():
                # It's a file - process by format
                self.add_files(inputs)
            else:
                # Assume it's text strings
                self.add_texts(inputs)
        else:
            # Assume it's text strings
            self.add_texts(inputs)
    
    def add_files(self, file_paths: List[Union[str, Path]], replace_if_exists: bool = True):
        """
        Add files - automatically detects format and processes accordingly.
        
        Args:
            file_paths: List of file paths (PDF, CSV, Excel, Word, PowerPoint, Text)
            replace_if_exists: If True, replace existing documents from the same files
        """
        pdf_files = []
        csv_files = []
        excel_files = []
        word_files = []
        powerpoint_files = []
        text_files = []
        
        for file_path in file_paths:
            file_path = Path(file_path)
            
            if not file_path.exists():
                print(f"⚠️  Skipping non-existent file: {file_path}")
                continue
            
            format_type = self._detect_format(file_path)
            
            if format_type == 'pdf':
                pdf_files.append(file_path)
            elif format_type == 'csv':
                csv_files.append(file_path)
            elif format_type == 'excel':
                excel_files.append(file_path)
            elif format_type == 'word':
                word_files.append(file_path)
            elif format_type == 'powerpoint':
                powerpoint_files.append(file_path)
            elif format_type == 'text':
                text_files.append(file_path)
            else:
                print(f"⚠️  Skipping unsupported format: {file_path}")
        
        # Process PDF files
        if pdf_files:
            self.add_pdf(pdf_files, replace_if_exists=replace_if_exists)
        
        # Process CSV files
        if csv_files:
            if PANDAS_AVAILABLE or CSV_EXCEL_LOADER_AVAILABLE:
                self.add_csv(csv_files, replace_if_exists=replace_if_exists)
            else:
                print("⚠️  CSV files skipped - CSV/Excel support not available. Install: pip install pandas")
        
        # Process Excel files
        if excel_files:
            if PANDAS_AVAILABLE or CSV_EXCEL_LOADER_AVAILABLE:
                self.add_excel(excel_files, replace_if_exists=replace_if_exists)
            else:
                print("⚠️  Excel files skipped - CSV/Excel support not available. Install: pip install pandas openpyxl")
        
        # Process Word files
        if word_files:
            if WORD_LOADER_AVAILABLE:
                self.add_word(word_files, replace_if_exists=replace_if_exists)
            else:
                print("⚠️  Word files skipped - Word support not available. Install: pip install python-docx")
        
        # Process PowerPoint files
        if powerpoint_files:
            if POWERPOINT_LOADER_AVAILABLE:
                self.add_powerpoint(powerpoint_files, replace_if_exists=replace_if_exists)
            else:
                print("⚠️  PowerPoint files skipped - PowerPoint support not available. Install: pip install python-pptx or pip install unstructured")
        
        # Process text files
        if text_files:
            self.add_text_file(text_files, replace_if_exists=replace_if_exists)
        
        # Re-initialize keyword retriever if needed
        if self.use_hybrid_search:
            self._init_keyword_retriever()
        
        # Ensure retriever is set if we have documents
        if not self.retriever and self.vectorstore:
            self.retriever = self.vectorstore.as_retriever(search_kwargs={"k": self.retrieval_k})
        
        print("✅ All files processed!")


def main():
    """Main demo."""
    print("=" * 60)
    print("Multi-Format RAG Example")
    print("=" * 60)
    
    print("\n1. Initializing...")
    # Use config file for all settings - modify settings.py to change settings
    try:
        from rag.settings import get_config
    except ImportError:
        # Fallback for direct script execution
        sys.path.insert(0, str(Path(__file__).parent))
        from settings import get_config
    config = get_config()
    rag = MultiFormatRAG(config=config)
    
    # Use delete_all=True to delete all ChromaDB collections from previous runs
    rag.delete_chromadb(delete_all=False)
    
    print("\n2. Adding mixed format documents...")
    # Example: Add text documents
    texts = [
        ""
    ]
    rag.add_texts(texts)
    
    # Example: Add files (would need actual file paths)
    files = [
         "/Users/mukesh/Downloads/Resume-Mukesh-Head-of-QA.pdf",
         "/Users/mukesh/Downloads/cards.xlsx"
    ]
    if files:
        rag.add_files(files)
    
    print("\n3. Asking questions...")
    questions = [
        "How many years of experience does Mukesh have?",
        "Give me steps to create a physical card",
    ]
    
    for q in questions:
        print("-" * 60)
        result = rag.query(q)
        print()
    
if __name__ == "__main__":
    main()
